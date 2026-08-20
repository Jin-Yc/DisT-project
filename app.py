"""Persisted product-readiness workflow for the DisT prototype."""

from __future__ import annotations

import copy
import json
import os
import sqlite3
import urllib.error
import urllib.request
from datetime import date
from io import BytesIO
from contextlib import closing
from pathlib import Path

from flask import Flask, jsonify, redirect, request, send_file, send_from_directory
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from pptx.util import Pt


ROOT = Path(__file__).resolve().parent
TEAM_ROLES = ("Dsci", "DA & RV", "Ops")
ROLE_FOCUS = {
    "Dsci": "方法论可行性、标准对齐、客户预期与痛点",
    "DA & RV": "数据覆盖、数据粒度、更新频率与招募要求",
    "Ops": "Scope、上线时间、成功 KPI 与交付质量",
}
ITERATION_PROJECTS = {
    "ri": {
        "id": "ri", "name": "RI", "type": "现有产品迭代", "stage": "方案优化", "readiness": "1 个待确认事项", "next": "确认新版访谈问题", "owner": "Alice",
        "objective": "优化访谈洞察的结构化输出，让销售团队更快识别高潜客户需求。",
        "iteration": "v1.2 · 访谈框架与洞察标签优化",
        "scope": "本轮聚焦访谈提纲、洞察标签和复盘模板；不改变现有客户分层模型。",
        "risk": "需要确认新版访谈问题是否覆盖关键购买决策。",
        "tasks": {
            "Dsci": {"title": "复核新版访谈框架", "due_date": "2026-08-13", "status": "待确认"},
            "DA & RV": {"title": "验证洞察标签覆盖率", "due_date": "2026-08-15", "status": "进行中"},
            "Ops": {"title": "更新客户访谈执行指引", "due_date": "2026-08-16", "status": "待开始"},
        },
    },
    "ecom": {
        "id": "ecom", "name": "Ecom", "type": "现有产品迭代", "stage": "开发验证", "readiness": "1 个交付风险", "next": "确认验收标准", "owner": "Wei",
        "objective": "提升电商经营看板的异常识别效率，并让运营团队能更快完成促销复盘。",
        "iteration": "v2.4 · 异常提醒与促销复盘优化",
        "scope": "本轮增加异常原因提示并调整促销复盘视图；不新增自定义报表能力。",
        "risk": "验收标准尚未与运营团队完成确认，可能影响本轮上线时间。",
        "tasks": {
            "Dsci": {"title": "校验异常提醒阈值", "due_date": "2026-08-12", "status": "进行中"},
            "DA & RV": {"title": "确认转化漏斗数据口径", "due_date": "2026-08-14", "status": "待确认"},
            "Ops": {"title": "确认 Ecom 验收标准", "due_date": "2026-08-11", "status": "待确认"},
        },
    },
}
def build_app(database_path: Path | None = None) -> Flask:
    app = Flask(__name__, static_folder="assets", static_url_path="/assets")
    db_path = database_path or ROOT / "instance" / "dist.db"
    db_path.parent.mkdir(parents=True, exist_ok=True)
    app.config["DATABASE"] = str(db_path)

    def connection() -> sqlite3.Connection:
        db = sqlite3.connect(app.config["DATABASE"])
        db.row_factory = sqlite3.Row
        return db

    def initialise() -> None:
        with closing(connection()) as db, db:
            db.execute("CREATE TABLE IF NOT EXISTS pl_project_state (project_id TEXT PRIMARY KEY, payload TEXT NOT NULL)")
            db.execute("DELETE FROM pl_project_state WHERE project_id = 'o2o' OR json_extract(payload, '$.name') = 'O2O'")
            db.execute("DROP TABLE IF EXISTS workflow_state")

    def read_pl_projects() -> list[dict]:
        with closing(connection()) as db, db:
            rows = db.execute("SELECT payload FROM pl_project_state ORDER BY rowid DESC").fetchall()
        return [normalise_pl_project(item) for item in (json.loads(row["payload"]) for row in rows) if item.get("id") != "o2o" and item.get("name") != "O2O"]

    def write_pl_project(state: dict) -> dict:
        state = normalise_pl_project(state)
        with closing(connection()) as db, db:
            db.execute("INSERT OR REPLACE INTO pl_project_state VALUES (?, ?)", (state["id"], json.dumps(state, ensure_ascii=False)))
        return state

    def normalise_pl_project(state: dict) -> dict:
        """Keep previously saved proposal payloads safe to consume as full records."""
        state.setdefault("context", {})
        state.setdefault("final_plan", {})
        state.setdefault("minutes", "")
        state.setdefault("minutes_analysis", "")
        state.setdefault("minutes_updates", [])
        state.setdefault("team_instructions", {
            role: f"评审 {state.get('name', '当前项目')} 的最终方案：重点确认 {ROLE_FOCUS[role]}。"
            for role in TEAM_ROLES
        })
        state.setdefault("issues", [])
        # Old payloads had one review task collection and jumped directly to
        # Leader Check.  Preserve those records as a completed final round,
        # while all newly created projects use the explicit two-round flow.
        legacy_tasks = state.setdefault("review_tasks", {})
        if "review_phase" not in state:
            state["review_phase"] = "final_complete" if state.get("stage") == "等待 PL 确认" else "initial_review"
        state.setdefault("report_version", 1)
        state.setdefault("meeting_started", False)
        state.setdefault("meeting_items_recorded", False)
        state.setdefault("first_review_tasks", copy.deepcopy(legacy_tasks))
        state.setdefault("final_review_tasks", copy.deepcopy(legacy_tasks) if state["review_phase"] == "final_complete" else {})
        for role in TEAM_ROLES:
            state["first_review_tasks"].setdefault(role, {"status": "pending", "conclusion": ""})
            state["final_review_tasks"].setdefault(role, {"status": "pending", "conclusion": ""})
        # review_tasks remains a compatible view for existing consumers.
        state["review_tasks"] = state["final_review_tasks"] if state["review_phase"] in {"final_review", "final_complete"} else state["first_review_tasks"]
        for issue in state["issues"]:
            issue.setdefault("resolution", "direct")
            if issue.get("status") == "awaiting_submitter":
                issue["resolution"] = "direct"
            if issue.get("status") == "meeting_required":
                issue["resolution"] = "meeting_required"
        checks = state.setdefault("leader_checks", {})
        for role in ("PL", *TEAM_ROLES):
            checks.setdefault(role, {"viewed": False, "confirmed": False, "note": ""})
        state.setdefault("confirmed", False)
        state.setdefault("stage", "团队评审中")
        state.setdefault("readiness", "等待团队评审")
        schedule = state.setdefault("schedule", {})
        schedule.setdefault("minutes", "")
        schedule.setdefault("analysis", "")
        schedule.setdefault("team_schedules", {role: "待排期会议分析" for role in TEAM_ROLES})
        schedule.setdefault("work_packages", [])
        schedule.setdefault("confirmed", False)
        schedule.setdefault("milestones", [])
        # Kept for backwards-compatible reads of historic records.  New
        # schedules distribute structured work packages instead.
        schedule.setdefault("tasks", {})
        return state

    def project_workflow_kind(project: dict) -> str:
        return "迭代项目" if project.get("context", {}).get("projectType") == "产品迭代" else "PL 新项目"

    def public_pl_project(state: dict, role: str | None = None) -> dict:
        project = copy.deepcopy(normalise_pl_project(state))
        project["open_issue_count"] = sum(item["status"] != "closed" for item in project["issues"])
        project["active_review_tasks"] = project["final_review_tasks"] if project["review_phase"] in {"final_review", "final_complete"} else project["first_review_tasks"]
        if role in TEAM_ROLES:
            project["role"] = role
            project["role_focus"] = ROLE_FOCUS[role]
            project["role_review_task"] = project["review_tasks"].get(role)
            project["my_issues"] = [item for item in project["issues"] if item["owner_role"] == role]
        return project

    def pl_pending_project(project: dict) -> dict:
        phase = project["review_phase"]
        first_review_ready = all(task["status"] == "submitted" for task in project["first_review_tasks"].values())
        final_review_ready = all(task["status"] == "submitted" for task in project["final_review_tasks"].values())
        leader_check_ready = all(check["confirmed"] for role, check in project["leader_checks"].items() if role in TEAM_ROLES)
        action = {
            "initial_review": "启动会议" if first_review_ready else "等待三方首次评审",
            "meeting": "更新报告并发起最终评审",
            "final_review": "完成最终评审" if final_review_ready else "查看最终评审进度",
            "final_complete": "确认项目" if leader_check_ready else "等待最终评审确认",
        }[phase]
        return {
            "kind": project_workflow_kind(project),
            "project_id": project["id"],
            "project": project["name"],
            "title": f"{project['name']}：{action}",
            "phase": phase,
            "action": action,
            "state": project["stage"],
            "status": project["readiness"],
        }

    def error(message: str, status: int = 409):
        return jsonify({"error": message}), status

    @app.before_request
    def retire_o2o_api():
        if request.path.startswith("/api/o2o"):
            return error("O2O 示例工作流已退役；请由 PL 新建项目开始。", 410)

    def positioning_fallback(context: dict, message: str = "") -> dict:
        name = context.get("productName") or "当前方案"
        return {
            "reply": f"已记录。围绕{name}，下一步请明确目标用户会在什么场景下遇到这一问题，以及首版必须验证的能力。",
            "summary": f"{name} 已具备初步方向，仍需收敛客户、场景与首版范围。",
            "suggestion": "先确认目标客户、核心痛点和最小可验证能力。",
            "assumption": "目标用户愿意改变现有工作方式来使用该产品。",
        }

    def positioning_assistant(context: dict, messages: list[dict]) -> dict:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            return positioning_fallback(context, messages[-1].get("content", "") if messages else "")
        prompt = f"""You are an expert product positioning assistant. Return strict JSON with reply, summary, suggestion, assumption.\nProduct: {context.get('productName', 'Untitled')}\nType: {context.get('projectType', 'New product')}\nBackground: {context.get('projectDesc', '')}\nReference: {context.get('projectRefer', '')}"""
        payload = {"model": "gpt-4o-mini", "messages": [{"role": "system", "content": prompt}, *messages], "temperature": 0.7, "max_tokens": 400}
        try:
            request_data = urllib.request.Request("https://api.openai.com/v1/chat/completions", data=json.dumps(payload).encode(), headers={"Content-Type": "application/json", "Authorization": f"Bearer {api_key}"}, method="POST")
            with urllib.request.urlopen(request_data, timeout=60) as response:
                parsed = json.loads(json.loads(response.read().decode())["choices"][0]["message"]["content"])
            return {key: parsed.get(key) or positioning_fallback(context)[key] for key in ("reply", "summary", "suggestion", "assumption")}
        except Exception:
            return positioning_fallback(context)

    def minutes_fallback(context: dict, text: str, current_plan: dict | None = None) -> dict:
        name = context.get("productName") or "当前方案"
        plan = copy.deepcopy(current_plan) if isinstance(current_plan, dict) else {
            "positioning": {"background": context.get("projectDesc") or "根据会议结论更新产品定位。", "target": "产品经理、业务负责人和研发评审人", "pain": "方案缺少清晰的决策依据", "value": "将讨论结论沉淀为可评审方案", "pitch": "从模糊想法到可执行方案"},
            "scope": {"inScope": [], "outScope": []}, "business": {"data": [], "metrics": [], "delivery": [], "roles": []},
            "risks": {"risks": [], "dependencies": [], "assumptions": [], "pending": []}, "actions": [],
        }
        excerpt = (text or "未提供文字内容").strip().replace("\n", " ")[:60]
        scope = plan.setdefault("scope", {"inScope": [], "outScope": []})
        business = plan.setdefault("business", {"data": [], "metrics": [], "delivery": [], "roles": []})
        risks = plan.setdefault("risks", {"risks": [], "dependencies": [], "assumptions": [], "pending": []})
        actions = plan.setdefault("actions", [])
        updates = []
        maturity = plan.setdefault("maturity", {})
        maturity_deltas = {}
        def append_once(items: list, value: str):
            if value not in items:
                items.append(value)
        def improve(key: str, keywords: tuple[str, ...], label: str):
            if not any(word in text for word in keywords):
                return
            previous = int(maturity.get(key, 5))
            current = min(10, previous + 1)
            maturity[key] = current
            if current > previous:
                maturity_deltas[key] = current - previous
                updates.append(f"成熟度：{label} {previous}/10 → {current}/10")
        if any(word in text for word in ("范围", "功能", "优先")):
            append_once(scope.setdefault("inScope", []), "会议确认的优先范围")
            updates.append("范围定义：已补充会议确认的优先范围")
        if any(word in text for word in ("风险", "依赖", "阻塞")):
            append_once(risks.setdefault("risks", []), "会议纪要中识别的待跟进风险")
            updates.append("风险与假设：已补充待跟进风险")
        if any(word in text for word in ("数据", "口径")):
            append_once(business.setdefault("data", []), "会议确认的数据口径与覆盖范围")
            updates.append("业务要求：已补充数据口径与覆盖范围")
        if any(word in text for word in ("KPI", "指标", "衡量")):
            append_once(business.setdefault("metrics", []), "会议确认的试点成功指标")
            updates.append("业务要求：已补充试点成功指标")
        improve("client", ("客户", "试点"), "客户理解")
        improve("pain", ("痛点", "异常"), "痛点清晰度")
        improve("positioning", ("定位", "价值"), "产品定位")
        improve("scope", ("范围", "首版", "优先"), "Scope 完整度")
        improve("execution", ("行动", "研发", "负责人"), "可执行性")
        improve("business", ("KPI", "指标", "数据"), "商业化准备度")
        append_once(actions, "跟进会议结论：" + excerpt)
        updates.append("下一步行动：已添加会议结论跟进项")
        plan.setdefault("positioning", {}).setdefault("background", context.get("projectDesc") or "根据会议结论更新产品定位。")
        plan["maturityDeltas"] = maturity_deltas
        plan["verdict"] = {"title": "会议结论已同步，待最终确认", "text": f"已从会议纪要提取结论并更新当前方案：{excerpt}"}
        return {"analysis": f"已模拟识别{name}会议纪要中的范围、风险与行动项，并同步到当前报告。", "updates": updates, "maturityDeltas": maturity_deltas, "updatedPlan": plan}

    def answer_plan_question(report: dict, question: str) -> str:
        query = question.lower()
        if any(word in query for word in ("范围", "scope")):
            scope = report.get("scope", {})
            return f"当前包含：{'；'.join(scope.get('inScope', [])) or '暂无'}。不包含：{'；'.join(scope.get('outScope', [])) or '暂无'}。"
        if any(word in query for word in ("结论", "建议")):
            verdict = report.get("verdict", {})
            return f"方案结论是“{verdict.get('title', '暂无结论')}”。{verdict.get('text', '')}"
        if any(word in query for word in ("风险", "假设", "依赖")):
            risks = report.get("risks", {})
            return f"主要风险：{'；'.join(risks.get('risks', [])) or '暂无'}。关键依赖：{'；'.join(risks.get('dependencies', [])) or '暂无'}。"
        positioning = report.get("positioning", {})
        return f"方案背景：{positioning.get('background', '暂无')}；目标客户：{positioning.get('target', '暂无')}；核心痛点：{positioning.get('pain', '暂无')}。"

    def answer_review_assistant_question(project: dict, role: str, question: str) -> str:
        answer = answer_plan_question(project.get("final_plan", {}), question)
        return f"从 {role} 的评审重点“{ROLE_FOCUS[role]}”来看，{answer} 如需形成正式意见，请将待确认事项提交为 Issue。"

    def ensure_team_role(role: str):
        if role not in TEAM_ROLES:
            return error("请选择有效的团队角色。", 400)
        return None

    initialise()

    @app.get("/")
    @app.get("/index.html")
    def index():
        return send_from_directory(ROOT, "index.html")

    @app.get("/workflow.html")
    def workflow():
        return redirect("/new-project.html")

    @app.get("/role-workflow.html")
    def role_workflow():
        return redirect("/index.html")

    @app.get("/role-review.html")
    def role_review():
        return send_from_directory(ROOT, "role-review.html")

    @app.get("/project-view.html")
    def project_view():
        return send_from_directory(ROOT, "project-view.html")

    @app.get("/new-project.html")
    def new_project():
        return send_from_directory(ROOT, "new-project.html")

    @app.get("/demo.html")
    def demo():
        return redirect("/index.html")

    @app.get("/api/overview")
    def overview():
        role = request.args.get("role", "PL")
        projects = [
            {key: value for key, value in project.items() if key in {"id", "name", "type", "stage", "readiness", "next"}}
            for project in ITERATION_PROJECTS.values()
        ]
        pl_projects = read_pl_projects()
        confirmed_pl_projects = [item for item in pl_projects if item.get("confirmed")]
        projects.extend({"id": item["id"], "name": item["name"], "type": "产品迭代" if project_workflow_kind(item) == "迭代项目" else "PL 新项目", "stage": item["stage"], "readiness": item["readiness"], "next": "查看项目概览", "pl_project": True} for item in confirmed_pl_projects)
        if role == "PL":
            pending_projects = [pl_pending_project(item) for item in pl_projects if not item.get("confirmed")]
            gantt = [{"project": item["name"], **package} for item in confirmed_pl_projects if item["schedule"].get("confirmed") for package in item["schedule"].get("work_packages", [])]
            return jsonify({"projects": projects, "role": role, "pending_projects": pending_projects, "gantt": gantt})
        if role not in TEAM_ROLES:
            return error("未识别的角色。", 400)
        tasks = [
            {"kind": "迭代", "title": project["tasks"][role]["title"], "status": project["tasks"][role]["status"], "due_date": project["tasks"][role]["due_date"], "project": project["name"], "project_id": project["id"]}
            for project in ITERATION_PROJECTS.values()
        ]
        for project in pl_projects:
            phase = project["review_phase"]
            task = (project["first_review_tasks"] if phase in {"initial_review", "meeting"} else project["final_review_tasks"]).get(role, {})
            actionable = (
                (phase == "initial_review" and task.get("status") == "pending")
                or (phase == "final_review" and task.get("status") == "pending")
            )
            if actionable:
                action = "提交首次评审（可确认无 Issue）" if phase == "initial_review" else "确认会议项已解决并提交最终结论"
                tasks.append({"kind": project_workflow_kind(project), "title": f"{project['name']}：{action}", "status": task["status"],
                              "phase": phase, "action": action, "state": project["stage"], "due_date": "待处理",
                              "project": project["name"], "project_id": project["id"], "pl_project": True, "review_task": True})
            if phase == "final_complete" and not project["leader_checks"][role]["confirmed"]:
                tasks.append({"kind": project_workflow_kind(project), "title": f"{project['name']}：完成 Leader Check", "status": "pending",
                              "phase": phase, "action": "确认最终方案", "state": project["stage"], "due_date": "待处理",
                              "project": project["name"], "project_id": project["id"], "pl_project": True, "review_task": True})
            if project.get("confirmed") and project["schedule"].get("confirmed"):
                for package in project["schedule"].get("work_packages", []):
                    if package.get("role") == role:
                        tasks.append({"kind": "迭代" if project_workflow_kind(project) == "迭代项目" else "新产品",
                                      "title": package["title"], "status": package.get("status", "待开始"),
                                      "due_date": package.get("due_date", "待排期"), "project": project["name"],
                                      "project_id": project["id"], "pl_project": True, "execution_task": True})
        gantt = [{"project": item["name"], **package} for item in confirmed_pl_projects if item["schedule"].get("confirmed") for package in item["schedule"].get("work_packages", []) if package.get("role") == role]
        return jsonify({"projects": projects, "role": role, "role_focus": ROLE_FOCUS[role], "tasks": tasks, "gantt": gantt})

    @app.post("/api/pl-projects")
    def create_pl_project():
        body = request.get_json(silent=True) or {}
        project_id = str(body.get("id", "")).strip()
        context = body.get("context") or {}
        if not project_id or not context.get("productName"):
            return error("项目编号和名称不能为空。", 400)
        state = {
            "id": project_id, "name": context["productName"], "context": copy.deepcopy(context),
            "final_plan": copy.deepcopy(body.get("finalPlan") or {}),
            "minutes": str(body.get("minutes", "")), "minutes_analysis": str(body.get("minutesAnalysis", "")),
            "minutes_updates": copy.deepcopy(body.get("minutesUpdates") or []),
            "stage": "首次团队评审中", "readiness": "等待三方首次评审", "confirmed": False,
            "review_phase": "initial_review", "report_version": 1, "meeting_started": False, "meeting_items_recorded": False,
            "team_instructions": {role: f"评审 {context['productName']} 的最终方案：重点确认 {ROLE_FOCUS[role]}。" for role in TEAM_ROLES},
            "issues": [], "first_review_tasks": {role: {"status": "pending", "conclusion": ""} for role in TEAM_ROLES},
            "final_review_tasks": {role: {"status": "pending", "conclusion": ""} for role in TEAM_ROLES},
            "leader_checks": {role: {"viewed": False, "confirmed": False, "note": ""} for role in ("PL", *TEAM_ROLES)},
        }
        return jsonify(public_pl_project(write_pl_project(state), "PL"))

    @app.delete("/api/pl-projects/unconfirmed")
    def revoke_unconfirmed_pl_projects():
        """Remove proposal-only projects and the team tasks derived from them."""
        body = request.get_json(silent=True) or {}
        if body.get("role") != "PL" and request.args.get("role") != "PL":
            return error("仅 PL 可以撤销未确认的新项目及其团队任务。", 403)
        try:
            with closing(connection()) as db, db:
                rows = db.execute("SELECT project_id, payload FROM pl_project_state").fetchall()
                project_ids = [
                    row["project_id"] for row in rows
                    if not bool(json.loads(row["payload"]).get("confirmed", False))
                ]
                if project_ids:
                    db.executemany(
                        "DELETE FROM pl_project_state WHERE project_id = ?",
                        [(project_id,) for project_id in project_ids],
                    )
        except (sqlite3.Error, json.JSONDecodeError) as exc:
            return error(f"撤销未确认项目失败：{exc}", 500)
        return jsonify({"deleted_count": len(project_ids)})

    @app.get("/api/pl-projects/<project_id>")
    def get_pl_project(project_id):
        role = request.args.get("role")
        if role and role not in TEAM_ROLES and role != "PL":
            return error("未识别的角色。", 400)
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        return jsonify(public_pl_project(project, role)) if project else error("未找到该新项目。", 404)

    @app.post("/api/pl-projects/<project_id>/issues")
    def create_pl_project_issue(project_id):
        body = request.get_json(silent=True) or {}
        role = str(body.get("role", ""))
        if role not in TEAM_ROLES:
            return error("仅团队角色可以提交 Issue。", 403)
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        if not project or project.get("review_phase") not in {"initial_review", "final_review"} or project.get("confirmed"):
            return error("仅团队评审阶段可以新增 Issue。")
        title, detail = str(body.get("title", "")).strip(), str(body.get("detail", "")).strip()
        if not title or not detail:
            return error("请填写问题和背景说明。", 400)
        project.setdefault("issues", []).append({"id": max((item["id"] for item in project["issues"]), default=0) + 1, "owner_role": role, "status": "open", "title": title, "detail": detail, "category": str(body.get("category", "Scope")), "priority": str(body.get("priority", "中")), "pl_response": "", "review_phase": project["review_phase"]})
        return jsonify(public_pl_project(write_pl_project(project), role))

    @app.post("/api/pl-projects/<project_id>/review-assistant")
    def review_project_assistant(project_id):
        body = request.get_json(silent=True) or {}
        role, question = str(body.get("role", "")), str(body.get("question", "")).strip()
        if role not in TEAM_ROLES:
            return error("仅团队角色可以使用项目问答助手。", 403)
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        if not project:
            return error("未找到该新项目。", 404)
        if project.get("confirmed") or project.get("review_phase") != "initial_review":
            return error("项目问答助手仅在首次团队评审阶段可用。", 409)
        if not question:
            return error("请填写想了解的项目问题。", 400)
        return jsonify({"reply": answer_review_assistant_question(project, role, question)})

    @app.post("/api/pl-projects/<project_id>/issues/<int:issue_id>/respond")
    def respond_pl_project_issue(project_id, issue_id):
        body = request.get_json(silent=True) or {}
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        issue = next((item for item in (project or {}).get("issues", []) if item["id"] == issue_id), None)
        response = str(body.get("response", "")).strip()
        if not issue or not response:
            return error("请填写 PL 处理说明。", 400)
        if project.get("review_phase") not in {"initial_review", "final_review"} or issue["status"] != "open":
            return error("该 Issue 当前无法回复。")
        action = body.get("action", "direct")
        if action not in {"direct", "meeting_required"}:
            return error("请选择“可直接解决”或“需会议解决”。", 400)
        issue["pl_response"], issue["resolution"] = response, action
        issue["status"] = "awaiting_submitter" if action == "direct" else "meeting_required"
        project["meeting_items_recorded"] = any(item["status"] == "meeting_required" for item in project["issues"])
        if action == "meeting_required" and project["review_phase"] == "final_review":
            project.update({"review_phase": "meeting", "stage": "会议进行中", "readiness": "等待会议纪要更新报告"})
        return jsonify(public_pl_project(write_pl_project(project)))

    @app.post("/api/pl-projects/<project_id>/issues/<int:issue_id>/confirm")
    def confirm_pl_project_issue(project_id, issue_id):
        body = request.get_json(silent=True) or {}
        role = str(body.get("role", ""))
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        issue = next((item for item in (project or {}).get("issues", []) if item["id"] == issue_id), None)
        if not issue or issue["owner_role"] != role:
            return error("仅提出 Issue 的角色可以确认关闭。", 403)
        if issue["status"] == "awaiting_submitter" and project.get("review_phase") in {"initial_review", "final_review"}:
            pass
        elif issue["status"] == "meeting_required" and project.get("review_phase") == "final_review":
            pass
        else:
            return error("直接回复须在首次评审确认；会议项须在最终评审确认已解决。")
        issue["status"] = "closed"
        return jsonify(public_pl_project(write_pl_project(project), role))

    @app.post("/api/pl-projects/<project_id>/reviews")
    def submit_pl_project_review(project_id):
        body = request.get_json(silent=True) or {}
        role, conclusion = str(body.get("role", "")), str(body.get("conclusion", "")).strip()
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        if role not in TEAM_ROLES or not project:
            return error("无法提交当前评审。", 400)
        task = project.get("first_review_tasks", {}).get(role)
        if project.get("review_phase") != "initial_review" or not task or task.get("status") != "pending":
            return error("当前评审任务不可提交。")
        if any(item["owner_role"] == role and item["status"] in {"open", "awaiting_submitter"} for item in project.get("issues", [])):
            return error("请先完成本角色可直接解决的 Issue。")
        if not conclusion:
            return error("请填写评审结论。", 400)
        project["first_review_tasks"][role] = {"status": "submitted", "conclusion": conclusion}
        return jsonify(public_pl_project(write_pl_project(project), role))

    @app.post("/api/pl-projects/<project_id>/meeting/start")
    def start_pl_project_meeting(project_id):
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        if not project:
            return error("未找到该新项目。", 404)
        if project.get("review_phase") != "initial_review":
            return error("当前不在可启动会议的首次评审阶段。")
        direct_open = any(item["status"] in {"open", "awaiting_submitter"} for item in project["issues"])
        reviews_ready = all(project["first_review_tasks"][role]["status"] == "submitted" for role in TEAM_ROLES)
        if direct_open or not reviews_ready:
            return error("需等待三方首次评审完成，并由提出者确认所有直接回复；会议项必须已记录。")
        project.update({"review_phase": "meeting", "meeting_started": True, "stage": "会议进行中", "readiness": "等待会议纪要更新报告"})
        return jsonify(public_pl_project(write_pl_project(project)))

    @app.post("/api/pl-projects/<project_id>/meeting")
    def apply_pl_project_meeting(project_id):
        body = request.get_json(silent=True) or {}
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        minutes = str(body.get("minutes", "")).strip()
        if not project:
            return error("未找到该新项目。", 404)
        if project.get("review_phase") != "meeting":
            return error("请先完成首次评审并启动会议。")
        if not minutes:
            return error("请粘贴或上传会议纪要。", 400)
        result = minutes_fallback(project["context"], minutes, project.get("final_plan"))
        project.update({"minutes": minutes, "minutes_analysis": result["analysis"], "minutes_updates": result["updates"],
                        "final_plan": result["updatedPlan"], "report_version": int(project.get("report_version", 1)) + 1,
                        "review_phase": "final_review", "stage": "最终评审中", "readiness": "等待三方最终结论"})
        project["final_review_tasks"] = {role: {"status": "pending", "conclusion": ""} for role in TEAM_ROLES}
        for role in TEAM_ROLES:
            project["leader_checks"][role] = {"viewed": False, "confirmed": False, "note": ""}
        return jsonify(public_pl_project(write_pl_project(project)))

    @app.post("/api/pl-projects/<project_id>/final-reviews")
    def submit_pl_project_final_review(project_id):
        body = request.get_json(silent=True) or {}
        role, conclusion = str(body.get("role", "")), str(body.get("conclusion", "")).strip()
        if role not in TEAM_ROLES:
            return error("仅团队角色可以提交最终结论。", 400)
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        if not project or project.get("review_phase") != "final_review":
            return error("当前不在最终评审阶段。")
        task = project["final_review_tasks"][role]
        if task["status"] != "pending":
            return error("本角色最终评审已提交。")
        if any(item["owner_role"] == role and item["status"] != "closed" for item in project["issues"]):
            return error("请先确认本角色会议项已解决。")
        if not conclusion:
            return error("请填写最终评审结论。", 400)
        project["final_review_tasks"][role] = {"status": "submitted", "conclusion": conclusion}
        project["leader_checks"][role].update({"viewed": True, "confirmed": True, "note": conclusion})
        return jsonify(public_pl_project(write_pl_project(project), role))

    @app.post("/api/pl-projects/<project_id>/final-reviews/complete")
    def complete_pl_project_final_reviews(project_id):
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        if not project:
            return error("未找到该新项目。", 404)
        if project.get("review_phase") != "final_review":
            return error("当前不在可完成的最终评审阶段。")
        if any(item["status"] != "closed" for item in project["issues"]) or not all(project["final_review_tasks"][role]["status"] == "submitted" for role in TEAM_ROLES):
            return error("需等待三方最终结论，并确认所有会议项已解决。")
        for role in TEAM_ROLES:
            project["final_review_tasks"][role]["status"] = "completed"
        project["leader_checks"]["PL"].update({"viewed": True, "confirmed": True, "note": "团队最终评审已完成"})
        project.update({"review_phase": "final_complete", "stage": "等待 PL 确认", "readiness": "最终评审完成"})
        return jsonify(public_pl_project(write_pl_project(project)))

    @app.post("/api/pl-projects/<project_id>/leader-checks")
    def update_pl_project_leader_check(project_id):
        body = request.get_json(silent=True) or {}
        role = str(body.get("role", ""))
        if role not in {"PL", *TEAM_ROLES}:
            return error("未识别的 Leader 角色。", 400)
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        if not project:
            return error("未找到该新项目。", 404)
        if project.get("review_phase") != "final_complete" or project.get("confirmed"):
            return error("最终评审完成后才能更新 Leader Check。")
        check = project["leader_checks"][role]
        check["viewed"] = True
        check["confirmed"] = bool(body.get("confirmed", True))
        check["note"] = str(body.get("note", "")).strip()
        return jsonify(public_pl_project(write_pl_project(project), role))

    @app.post("/api/pl-projects/<project_id>/confirm")
    def confirm_pl_project(project_id):
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        if not project:
            return error("未找到当前新项目，请先发起团队评审。", 404)
        ready = (
            project.get("review_phase") == "final_complete"
            and all(project["final_review_tasks"].get(role, {}).get("status") == "completed" for role in TEAM_ROLES)
            and not any(item.get("status") != "closed" for item in project.get("issues", []))
            and all(project["leader_checks"].get(role, {}).get("confirmed") for role in TEAM_ROLES)
        )
        if not ready:
            return error("确认项目需处于等待 PL 确认阶段，三方任务完成、Issue 全部关闭且所有 Leader Check 已确认。")
        project.update({"confirmed": True, "stage": "项目已确认", "readiness": "可进入项目执行"})
        return jsonify(public_pl_project(write_pl_project(project), "PL"))

    @app.post("/api/pl-projects/<project_id>/schedule")
    def update_pl_project_schedule(project_id):
        body = request.get_json(silent=True) or {}
        if body.get("role") != "PL":
            return error("仅 PL 可以维护项目排期。", 403)
        project = next((item for item in read_pl_projects() if item["id"] == project_id), None)
        if not project or not project.get("confirmed"):
            return error("仅已确认项目可以维护排期。", 409)
        milestones = body.get("milestones", [])
        team_schedules, work_packages = body.get("team_schedules", {}), body.get("work_packages", [])
        if not isinstance(milestones, list) or not isinstance(team_schedules, dict) or not isinstance(work_packages, list):
            return error("排期数据格式无效。", 400)
        cleaned_milestones = []
        for item in milestones:
            if not isinstance(item, dict):
                return error("排期节点格式无效。", 400)
            title, due_date = str(item.get("title", "")).strip(), str(item.get("due_date", "")).strip()
            if title or due_date:
                if not title or not due_date:
                    return error("每个关键节点必须包含名称和日期。", 400)
                try:
                    date.fromisoformat(due_date)
                except ValueError:
                    return error("关键节点日期必须为有效日期。", 400)
                cleaned_milestones.append({"title": title, "due_date": due_date})
        cleaned_schedules = {}
        for role in TEAM_ROLES:
            cleaned_schedules[role] = str(team_schedules.get(role, "")).strip()
        if not work_packages or any(not isinstance(item, dict) for item in work_packages):
            return error("请至少确认一个可分发的团队工作包。", 400)
        cleaned_packages = [{key: str(item.get(key, "")).strip() for key in ("role", "title", "start_date", "due_date", "dependency")} | {"status": str(item.get("status", "待开始")).strip() or "待开始"} for item in work_packages]
        if any(item["role"] not in TEAM_ROLES or not all(item[key] for key in ("title", "start_date", "due_date")) for item in cleaned_packages):
            return error("每个工作包必须包含团队、任务、开始日和截止日。", 400)
        try:
            if any(date.fromisoformat(item["start_date"]) > date.fromisoformat(item["due_date"]) for item in cleaned_packages):
                return error("工作包的开始日期不能晚于截止日期。", 400)
        except ValueError:
            return error("工作包日期必须为有效日期。", 400)
        project["schedule"].update({"minutes": str(body.get("minutes", "")).strip(), "milestones": cleaned_milestones, "team_schedules": cleaned_schedules, "work_packages": cleaned_packages, "confirmed": True})
        return jsonify(public_pl_project(write_pl_project(project), "PL"))

    @app.post("/api/pl-projects/<project_id>/schedule/analyze")
    def analyze_pl_project_schedule(project_id):
        project = next((item for item in read_pl_projects() if item["id"] == project_id and item.get("confirmed")), None)
        if not project:
            return error("仅已确认项目可以分析排期纪要。", 409)
        minutes = request.form.get("minutes", "").strip()
        uploaded = request.files.get("file")
        if uploaded and uploaded.filename:
            suffix = Path(uploaded.filename).suffix.lower()
            raw = uploaded.read()
            if suffix in {".txt", ".md"}:
                minutes = raw.decode("utf-8", errors="ignore") or minutes
            elif suffix == ".docx":
                minutes = "\n".join(paragraph.text for paragraph in Document(BytesIO(raw)).paragraphs if paragraph.text.strip()) or minutes
            elif suffix == ".pdf":
                minutes = "\n".join(page.extract_text() or "" for page in PdfReader(BytesIO(raw)).pages) or minutes
            else:
                return error("仅支持 TXT、Markdown、Word DOCX 和 PDF 纪要。", 400)
        if not minutes:
            return error("请粘贴或上传排期会议纪要。", 400)
        name = project["name"]
        suggestions = {role: f"围绕{ROLE_FOCUS[role]}确认 {name} 的排期依赖与交付顺序。" for role in TEAM_ROLES}
        milestones = [{"title": "范围与依赖确认", "due_date": "2026-08-21"}, {"title": "试点/交付验收", "due_date": "2026-08-29"}]
        packages = [{"role": "Dsci", "title": "方法与标准复核", "start_date": "2026-08-20", "due_date": "2026-08-23", "dependency": "范围确认", "status": "待开始"}, {"role": "DA & RV", "title": "数据口径与样本核对", "start_date": "2026-08-21", "due_date": "2026-08-26", "dependency": "方法复核", "status": "待开始"}, {"role": "Ops", "title": "上线窗口与验收准备", "start_date": "2026-08-26", "due_date": "2026-08-29", "dependency": "数据核对", "status": "待开始"}]
        return jsonify({"minutes": minutes, "analysis": f"已从排期纪要提取 {name} 的协作重点与日期草案；请由 PL 确认后分发。", "milestones": milestones, "team_schedules": suggestions, "work_packages": packages})

    @app.get("/api/pl-projects/<project_id>/export/pptx")
    def export_pl_project_pptx(project_id):
        project = next((item for item in read_pl_projects() if item["id"] == project_id and item.get("confirmed")), None)
        if not project:
            return error("仅已确认项目可以导出 PPT。", 409)
        plan, schedule = project.get("final_plan", {}), project.get("schedule", {})
        deck = Presentation()
        def add_slide(title: str, lines: list[str]):
            slide = deck.slides.add_slide(deck.slide_layouts[1]); slide.shapes.title.text = title
            frame = slide.placeholders[1].text_frame; frame.clear()
            for index, line in enumerate(lines):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph(); paragraph.text = line; paragraph.font.size = Pt(20)
        add_slide(project["name"], [project["context"].get("projectType", "产品方案"), project["context"].get("iterationVersion", "已确认方案")])
        positioning = plan.get("positioning", {}); add_slide("产品定位", [f"背景：{positioning.get('background', '待确认')}", f"目标客户：{positioning.get('target', '待确认')}", f"价值主张：{positioning.get('value', '待确认')}"])
        scope = plan.get("scope", {}); add_slide("范围定义", ["包含：" + "；".join(scope.get("inScope", [])), "暂不包含：" + "；".join(scope.get("outScope", []))])
        business = plan.get("business", {}); add_slide("业务要求", ["数据：" + "；".join(business.get("data", [])), "指标：" + "；".join(business.get("metrics", [])), "交付：" + "；".join(business.get("delivery", []))])
        risks = plan.get("risks", {}); add_slide("风险与依赖", ["风险：" + "；".join(risks.get("risks", [])), "依赖：" + "；".join(risks.get("dependencies", []))])
        add_slide("行动项与排期", [
            "节点：" + "；".join(f"{item.get('title', '')}（{item.get('due_date', '待定')}）" for item in schedule.get("milestones", [])),
            *[f"{item.get('role', '团队')}：{item.get('title', '待排期')}（{item.get('start_date', '待定')} 至 {item.get('due_date', '待定')}）" for item in schedule.get("work_packages", [])],
        ])
        output = BytesIO(); deck.save(output); output.seek(0)
        return send_file(output, as_attachment=True, download_name=f"{project['name']}-项目简报.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")

    @app.post("/api/positioning-assistant")
    def positioning_assistant_api():
        body = request.get_json(silent=True) or {}
        return jsonify(positioning_assistant(body.get("context", {}), body.get("messages", [])))

    @app.post("/api/plan-confirmation")
    def plan_confirmation_api():
        body = request.get_json(silent=True) or {}
        return jsonify({"answer": answer_plan_question(body.get("report", {}), body.get("question", ""))})

    @app.post("/api/meeting-minutes")
    def meeting_minutes_api():
        try:
            context = json.loads(request.form.get("context", "{}"))
        except json.JSONDecodeError:
            context = {}
        uploaded = request.files.get("file")
        text = request.form.get("content", "")
        try:
            current_plan = json.loads(request.form.get("currentPlan", "{}"))
        except json.JSONDecodeError:
            current_plan = {}
        if uploaded and uploaded.filename:
            text = uploaded.read().decode("utf-8", errors="ignore")
        return jsonify(minutes_fallback(context, text, current_plan))

    @app.get("/api/projects/<project_id>")
    def get_iteration_project(project_id: str):
        project = ITERATION_PROJECTS.get(project_id)
        role = request.args.get("role", "PL")
        if role not in {"PL", *TEAM_ROLES}:
            return error("未识别的角色。", 400)
        if not project:
            pl_project = next((item for item in read_pl_projects() if item["id"] == project_id and item.get("confirmed")), None)
            if not pl_project:
                return error("未找到该项目。", 404)
            plan = pl_project.get("final_plan", {})
            scope = plan.get("scope", {})
            risks = plan.get("risks", {})
            result = {
                "id": pl_project["id"], "name": pl_project["name"], "type": pl_project["context"].get("projectType", "新产品 Launch"),
                "stage": pl_project["stage"], "readiness": pl_project["readiness"], "objective": pl_project["context"].get("projectDesc", ""),
                "iteration": pl_project["context"].get("iterationVersion") or "已确认方案", "scope": "；".join(scope.get("inScope", [])) or "待补充范围",
                "risk": "；".join(risks.get("risks", [])) or "暂无已识别风险",
                "schedule": pl_project["schedule"], "pl_project": True, "context": pl_project["context"], "final_plan": plan,
            }
            result["role"] = role
            result["my_work_packages"] = [] if role == "PL" else [
                package for package in pl_project["schedule"].get("work_packages", []) if package.get("role") == role
            ]
            return jsonify(result)
        result = copy.deepcopy(project)
        result["role"] = role
        result["my_task"] = None if role == "PL" else result["tasks"][role]
        return jsonify(result)

    return app


app = build_app()


if __name__ == "__main__":
    app.run(host="127.0.0.1", debug=False, port=int(os.environ.get("PORT", "8765")))
